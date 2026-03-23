import asyncio
import hashlib
import io
import ipaddress
import json
import logging
import os
import re
import socket
import urllib.parse
from pathlib import Path
from typing import Any

logger = logging.getLogger("Axon.Loaders")

_MAX_FILE_BYTES = 100 * 1024 * 1024  # 100 MB


def _stable_file_id(path: str, kind: str = "file") -> str:
    """SHA-256 based stable ID from absolute path. Prevents basename collisions."""
    abspath = os.path.normcase(os.path.normpath(os.path.abspath(path)))
    return kind + "_" + hashlib.sha256(abspath.encode()).hexdigest()[:24]


def _check_file_size(path: str) -> None:
    """Raise ValueError if the file exceeds _MAX_FILE_BYTES."""
    size = os.path.getsize(path)
    if size > _MAX_FILE_BYTES:
        raise ValueError(
            f"File '{path}' is {size / (1024 * 1024):.1f} MB, " f"which exceeds the 100 MB limit."
        )


def _rewrite_github_url(url: str) -> str:
    """Rewrite GitHub web URLs to their raw/fetchable equivalents.

    Supported rewrites:
    - ``github.com/<owner>/<repo>/blob/<ref>/<path>``
      → ``raw.githubusercontent.com/<owner>/<repo>/<ref>/<path>``
    - ``gist.github.com/<user>/<gist_id>``
      → ``gist.githubusercontent.com/<user>/<gist_id>/raw``

    Raises ``ValueError`` for GitHub tree (directory listing) URLs because they
    cannot be fetched as a single document.  Use a specific file URL
    (``/blob/...``) or clone the repo and use ``ingest_path`` instead.

    All other URLs are returned unchanged.
    """
    # github.com blob URL → raw URL
    blob_match = re.match(
        r"^https?://github\.com/([^/]+)/([^/]+)/blob/([^/]+)/(.+)$", url, re.IGNORECASE
    )
    if blob_match:
        owner, repo, ref, path = blob_match.groups()
        return f"https://raw.githubusercontent.com/{owner}/{repo}/{ref}/{path}"

    # github.com tree URL → helpful error (directory listings can't be one doc)
    if re.match(r"^https?://github\.com/[^/]+/[^/]+/tree/", url, re.IGNORECASE):
        raise ValueError(
            "GitHub directory (tree) URLs cannot be fetched as a single document. "
            "Provide a specific file URL ending in /blob/<ref>/<file>, or clone "
            "the repository locally and use ingest_path instead."
        )

    # gist.github.com URL → raw gist URL
    gist_match = re.match(r"^https?://gist\.github\.com/([^/]+)/([a-f0-9]+)/?$", url, re.IGNORECASE)
    if gist_match:
        user, gist_id = gist_match.groups()
        return f"https://gist.githubusercontent.com/{user}/{gist_id}/raw"

    return url


def _extract_html_text(html: str) -> str:
    """Extract visible text from an HTML string using stdlib html.parser."""
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self._texts = []
            self._content_skip_tags = {"script", "style"}
            self._current_skip = False

        def handle_starttag(self, tag, attrs):
            if tag.lower() in self._content_skip_tags:
                self._current_skip = True

        def handle_endtag(self, tag):
            if tag.lower() in self._content_skip_tags:
                self._current_skip = False

        def handle_data(self, data):
            if not self._current_skip:
                stripped = data.strip()
                if stripped:
                    self._texts.append(stripped)

        def get_text(self) -> str:
            return " ".join(self._texts)

    parser = _TextExtractor()
    parser.feed(html)
    return parser.get_text()


class BaseLoader:
    """Base class for document loaders."""

    def load(self, path: str) -> list[dict[str, Any]]:
        raise NotImplementedError

    async def aload(self, path: str) -> list[dict[str, Any]]:
        """Async version of load."""
        return await asyncio.to_thread(self.load, path)


class FlexibleTableLoader(BaseLoader):
    """
    Robust loader for tabular data (CSV, TSV, etc.) that handles:
    1. Automatic delimiter detection (sniffing).
    2. Ragged/jagged rows (varying column counts).
    3. Missing or present headers.
    """

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        with open(path, encoding="utf-8", errors="ignore") as f:
            content = f.read()
        return self.load_text(content, source=path)

    def load_text(self, text: str, source: str = "raw_text") -> list[dict[str, Any]]:
        """Process raw text as a table."""
        import csv
        import io

        from axon.splitters import TableSplitter

        sample = text[:8192]

        # Detect delimiter
        try:
            # Priority 1: Sniffer
            dialect = csv.Sniffer().sniff(sample, delimiters=",\t|;")
            delimiter = dialect.delimiter
        except Exception:
            # Priority 2: Simple heuristics
            counts = {"\t": sample.count("\t"), ",": sample.count(","), "|": sample.count("|")}
            delimiter = max(counts, key=lambda k: counts[k])
            # If everything is 0, default to comma
            if counts[delimiter] == 0:
                delimiter = ","

        # Detect header
        try:
            has_header = csv.Sniffer().has_header(sample)
        except Exception:
            has_header = True

        f = io.StringIO(text)
        reader = csv.reader(f, delimiter=delimiter)
        all_rows = [r for r in reader if r]  # Skip empty rows

        if not all_rows:
            return []

        # Find the max column count across ALL rows
        max_cols = max(len(row) for row in all_rows)

        # Determine headers and data
        if has_header:
            base_headers = all_rows[0]
            data_rows = all_rows[1:]
        else:
            base_headers = [f"Col_{i}" for i in range(len(all_rows[0]))]
            data_rows = all_rows

        # Ensure header count matches max_cols
        headers = list(base_headers)
        while len(headers) < max_cols:
            headers.append(f"Extra_Col_{len(headers)}")
        # If actual headers were shorter than row 0, handle that
        headers = headers[:max_cols]

        final_rows = []
        for row in data_rows:
            row_dict = {}
            for i, val in enumerate(row):
                if i < len(headers):
                    row_dict[headers[i]] = val.strip()
            final_rows.append(row_dict)

        table_name = os.path.basename(source)
        splitter = TableSplitter(table_name=table_name)
        chunks = splitter.transform_rows(final_rows, headers)

        # Include table_name in the hash so that multiple tables from the same source
        # file (e.g. multiple sheets) get distinct IDs and do not collide.
        source_key = f"{source}::{table_name}"
        stable_base = (
            _stable_file_id(source, "table")
            if not source.startswith(("http://", "https://"))
            and (os.path.isabs(source) or os.sep in source)
            else ("table_" + hashlib.sha256(source_key.encode()).hexdigest()[:24])
        )
        documents = []
        for i, text_chunk in enumerate(chunks):
            documents.append(
                {
                    "id": f"{stable_base}_row_{i}",
                    "text": text_chunk,
                    "metadata": {"source": source, "type": "table", "row": i},
                }
            )
        return documents


class SmartTextLoader(BaseLoader):
    """
    Delegates to FlexibleTableLoader if the file looks structured,
    otherwise uses TextLoader.
    """

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        with open(path, encoding="utf-8", errors="ignore") as f:
            content = f.read()
        return self.load_text(content, source=path)

    def load_text(self, text: str, source: str = "raw_text") -> list[dict[str, Any]]:
        """Intelligently detect if text is a table or prose and load accordingly."""
        # Check if it looks like a table (multiple tabs or commas per line)
        sample = text[:2048]
        lines = sample.strip().split("\n")
        tab_count = sample.count("\t")
        comma_count = sample.count(",")

        is_likely_table = False
        if len(lines) > 1:
            avg_tabs = tab_count / len(lines)
            avg_commas = comma_count / len(lines)
            if avg_tabs > 1.5 or avg_commas > 2.0:
                is_likely_table = True

        if is_likely_table:
            return FlexibleTableLoader().load_text(text, source=source)
        else:
            # Fallback to standard document dictionary format.
            # Use the full source as the id (already unique when source is a
            # UUID doc_id); os.path.basename strips directory context and can
            # produce duplicate ids for non-path sources on Windows.
            return [
                {
                    "id": source,
                    "text": text,
                    "metadata": {"source": source, "type": "text"},
                }
            ]


class TextLoader(BaseLoader):
    """Loader for plain text files."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        with open(path, encoding="utf-8") as f:
            content = f.read()
        return [
            {
                "id": _stable_file_id(path),
                "text": content,
                "metadata": {"source": path, "type": "text"},
            }
        ]


class CodeFileLoader(BaseLoader):
    """Loader for source code files.

    Sets ``type="code"`` in metadata so the ``DirectoryLoader`` path-enrichment
    step skips prepending a ``[File Path: ...]`` prefix into the raw text.
    Keeping the prefix out of code text is required so that AST-based chunkers
    (e.g. ``CodeAwareSplitter._split_python_ast``) receive valid syntax.
    The file path is preserved in ``metadata["source"]``.
    """

    def load(self, path: str) -> list[dict[str, Any]]:
        import hashlib as _hl

        _check_file_size(path)
        with open(path, encoding="utf-8", errors="ignore") as f:
            content = f.read()
        stable_id = "code_" + _hl.sha256(os.path.abspath(path).encode("utf-8")).hexdigest()[:24]
        return [
            {
                "id": stable_id,
                "text": content,
                "metadata": {"source": path, "type": "code"},
            }
        ]


class TSVLoader(BaseLoader):
    """Loader for tab-delimited files. Each row becomes an enriched document chunk."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        import csv

        from axon.splitters import TableSplitter

        rows = []
        headers: list[str] = []
        with open(path, encoding="utf-8", newline="") as f:
            reader = csv.DictReader(f, delimiter="\t")
            headers = list(reader.fieldnames or [])
            for row in reader:
                rows.append(row)

        table_name = os.path.basename(path)
        splitter = TableSplitter(table_name=table_name)
        chunks = splitter.transform_rows(rows, headers)

        stable_base = _stable_file_id(path, "tsv")
        documents = []
        for i, text in enumerate(chunks):
            documents.append(
                {
                    "id": f"{stable_base}_row_{i}",
                    "text": text,
                    "metadata": {"source": path, "type": "tsv", "row": i},
                }
            )
        return documents


class JSONLoader(BaseLoader):
    """Loader for JSON files."""

    def _sanitize_metadata(self, metadata: dict[str, Any]) -> dict[str, Any]:
        """Ensure metadata values are scalar (str, int, float, bool)."""
        sanitized = {}
        for k, v in metadata.items():
            if isinstance(v, dict | list):
                sanitized[k] = json.dumps(v)
            else:
                sanitized[k] = v
        return sanitized

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        with open(path, encoding="utf-8") as f:
            try:
                data = json.load(f)
            except json.JSONDecodeError as exc:
                logger.warning(f"Skipping malformed JSON file {path}: {exc}")
                return []

        if isinstance(data, list):
            documents = []
            for i, item in enumerate(data):
                text = item.get("text", item.get("content", json.dumps(item)))
                metadata = {k: v for k, v in item.items() if k not in ["text", "content"]}
                metadata = self._sanitize_metadata(metadata)
                metadata.update({"source": path, "type": "json", "index": i})
                documents.append(
                    {
                        "id": f"{_stable_file_id(path, 'json')}_{i}",
                        "text": text,
                        "metadata": metadata,
                    }
                )
            return documents
        else:
            text = data.get("text", data.get("content", json.dumps(data)))
            metadata = {k: v for k, v in data.items() if k not in ["text", "content"]}
            metadata = self._sanitize_metadata(metadata)
            metadata.update({"source": path, "type": "json"})
            return [{"id": _stable_file_id(path, "json"), "text": text, "metadata": metadata}]


class CSVLoader(BaseLoader):
    """Loader for CSV files. Each row becomes an enriched document chunk."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        import csv

        from axon.splitters import TableSplitter

        rows = []
        headers: list[str] = []
        with open(path, encoding="utf-8", newline="") as f:
            reader = csv.DictReader(f)
            headers = list(reader.fieldnames or [])
            for row in reader:
                rows.append(row)

        table_name = os.path.basename(path)
        splitter = TableSplitter(table_name=table_name)
        chunks = splitter.transform_rows(rows, headers)

        stable_base = _stable_file_id(path, "csv")
        documents = []
        for i, text in enumerate(chunks):
            documents.append(
                {
                    "id": f"{stable_base}_row_{i}",
                    "text": text,
                    "metadata": {"source": path, "type": "csv", "row": i},
                }
            )
        return documents


class HTMLLoader(BaseLoader):
    """Loader for HTML files. Extracts visible text content."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        with open(path, encoding="utf-8", errors="ignore") as f:
            content = f.read()
        text = _extract_html_text(content)
        return [
            {
                "id": _stable_file_id(path, "html"),
                "text": text,
                "metadata": {"source": path, "type": "html"},
            }
        ]


class URLLoader(BaseLoader):
    """Loader for HTTP/HTTPS URLs with SSRF mitigations.

    Fetches the URL using httpx (already in project deps), strips HTML if the
    response is text/html, and returns a single document dict.

    Security:
    - Only http/https schemes are accepted; all others raise ValueError.
    - Resolved IP addresses are checked against blocked private/internal ranges
      to prevent Server-Side Request Forgery (SSRF).
    - Redirects are capped at 5 hops.
    - Response content type must be text/*; binary responses are rejected.
    - Content size is capped at _MAX_FILE_BYTES (100 MB).
    """

    _BLOCKED_NETWORKS: list[ipaddress.IPv4Network | ipaddress.IPv6Network] = [
        ipaddress.ip_network("127.0.0.0/8"),  # loopback
        ipaddress.ip_network("10.0.0.0/8"),  # RFC 1918 private
        ipaddress.ip_network("172.16.0.0/12"),  # RFC 1918 private
        ipaddress.ip_network("192.168.0.0/16"),  # RFC 1918 private
        ipaddress.ip_network("169.254.0.0/16"),  # link-local / cloud metadata (AWS/GCP/Azure)
        ipaddress.ip_network("::1/128"),  # IPv6 loopback
        ipaddress.ip_network("fc00::/7"),  # IPv6 unique local
        ipaddress.ip_network("fe80::/10"),  # IPv6 link-local
    ]

    def _check_ssrf(self, url: str) -> None:
        """Raise ValueError if the URL targets a private/internal address."""
        parsed = urllib.parse.urlparse(url)
        if parsed.scheme not in ("http", "https"):
            raise ValueError(
                f"Scheme '{parsed.scheme}' is not allowed. Only http and https are permitted."
            )
        hostname = parsed.hostname
        if not hostname:
            raise ValueError("No hostname found in URL.")
        try:
            addrinfos = socket.getaddrinfo(hostname, None)
        except socket.gaierror as exc:
            raise ValueError(f"Could not resolve hostname '{hostname}': {exc}")
        for addrinfo in addrinfos:
            ip_str = addrinfo[4][0]
            try:
                ip = ipaddress.ip_address(ip_str)
            except ValueError:
                continue
            for blocked in self._BLOCKED_NETWORKS:
                if ip in blocked:
                    raise ValueError(
                        f"Host '{hostname}' resolves to a blocked private/internal address ({ip})."
                    )

    def load(self, url: str) -> list[dict[str, Any]]:
        """Fetch *url* and return its text content as a single document.

        GitHub ``/blob/`` URLs are automatically rewritten to
        ``raw.githubusercontent.com`` before any network contact is made.
        GitHub Gist URLs are rewritten to their raw form.
        GitHub tree (directory) URLs raise ``ValueError`` immediately.
        """
        import httpx

        url = _rewrite_github_url(url)  # Gap-5: rewrite before SSRF check
        self._check_ssrf(url)
        try:
            with httpx.Client(timeout=30.0, follow_redirects=True, max_redirects=5) as client:
                resp = client.get(url)
        except httpx.TimeoutException:
            raise ValueError(f"Request to '{url}' timed out.")
        except httpx.TooManyRedirects:
            raise ValueError(f"Too many redirects fetching '{url}'.")
        except httpx.RequestError as exc:
            raise ValueError(f"Failed to fetch '{url}': {exc}")

        if resp.status_code != 200:
            raise ValueError(f"'{url}' returned HTTP {resp.status_code}.")

        content_type = resp.headers.get("content-type", "")
        if not content_type.startswith("text/"):
            raise ValueError(
                f"Non-text content type '{content_type}' for '{url}'. Only text/* is accepted."
            )

        raw = resp.text
        if len(raw.encode("utf-8")) > _MAX_FILE_BYTES:
            raise ValueError(
                f"URL content exceeds the {_MAX_FILE_BYTES // (1024 * 1024)} MB limit."
            )

        text = _extract_html_text(raw) if "html" in content_type else raw
        import hashlib as _hl

        stable_id = "url_" + _hl.sha256(url.encode("utf-8")).hexdigest()[:32]
        return [
            {
                "id": stable_id,
                "text": text,
                "metadata": {"source": url, "type": "url"},
            }
        ]


class DOCXLoader(BaseLoader):
    """Loader for DOCX files using python-docx."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        try:
            from docx import Document
        except ImportError:
            logger.error("python-docx not installed. Install with: pip install python-docx")
            return []

        doc = Document(path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        text = "\n\n".join(paragraphs)

        return [
            {
                "id": _stable_file_id(path, "docx"),
                "text": text,
                "metadata": {"source": path, "type": "docx"},
            }
        ]


class PPTXLoader(BaseLoader):
    """Loader for PPTX files using python-pptx."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return []

        prs = Presentation(path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text.strip())
        text = "\n\n".join([t for t in text_runs if t])

        return [
            {
                "id": _stable_file_id(path, "pptx"),
                "text": text,
                "metadata": {"source": path, "type": "pptx"},
            }
        ]


class ImageLoader(BaseLoader):
    """
    Loader for raster images (BMP, PNG, JPG, TIF/TIFF, PGM) using Ollama VLM for captioning.

    Images are converted to PNG bytes via Pillow before being sent to the VLM so that
    all formats — including exotic ones like PGM — are handled uniformly.
    """

    def __init__(self, ollama_model: str = "llava"):
        self.ollama_model = ollama_model
        try:
            from PIL import Image

            self._pil: Any = Image
        except ImportError:
            self._pil = None
            logger.warning("Pillow not installed. Image loading will be skipped.")
        try:
            import ollama

            self.ollama: Any = ollama
        except ImportError:
            self.ollama = None
            logger.warning("ollama package not installed. Image loading will be skipped.")

    def load(self, path: str) -> list[dict[str, Any]]:
        if self._pil is None or self.ollama is None:
            return []

        logger.info(f"🖼️ Processing image: {path} with {self.ollama_model}...")

        try:
            # Normalize to PNG bytes via Pillow for maximum VLM compatibility
            img = self._pil.open(path).convert("RGB")
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            image_data = buf.getvalue()

            fmt = Path(path).suffix.lstrip(".").lower()

            # Call Ollama VLM
            response = self.ollama.generate(
                model=self.ollama_model,
                prompt="Describe this image in detail. Mention any text, objects, people, or patterns visible.",
                images=[image_data],
            )

            description = response["response"]

            return [
                {
                    "id": _stable_file_id(path, "image"),
                    "text": f"Image Description: {description}",
                    "metadata": {
                        "source": path,
                        "type": "image",
                        "format": fmt,
                        "model": self.ollama_model,
                    },
                }
            ]
        except Exception as e:
            logger.error(f"Error processing image {path}: {e}")
            return []


# Backward-compatible alias kept for existing code that imports BMPLoader directly.
BMPLoader = ImageLoader


class PDFLoader(BaseLoader):
    """Loader for PDF files. Extracts text page-by-page using PyMuPDF (fitz) with pypdf fallback."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(path)
            total = len(doc)
            documents = []
            for page_num in range(total):
                page = doc[page_num]
                text = page.get_text()
                documents.append(
                    {
                        "id": f"{_stable_file_id(path, 'pdf')}_page_{page_num}",
                        "text": text,
                        "metadata": {
                            "source": path,
                            "type": "pdf",
                            "page": page_num,
                            "total_pages": total,
                        },
                    }
                )
            doc.close()
            return documents
        except ImportError:
            pass

        try:
            import pypdf

            reader = pypdf.PdfReader(path)
            total = len(reader.pages)
            documents = []
            for page_num in range(total):
                text = reader.pages[page_num].extract_text() or ""
                documents.append(
                    {
                        "id": f"{_stable_file_id(path, 'pdf')}_page_{page_num}",
                        "text": text,
                        "metadata": {
                            "source": path,
                            "type": "pdf",
                            "page": page_num,
                            "total_pages": total,
                        },
                    }
                )
            return documents
        except ImportError:
            pass

        logger.error(
            "Neither PyMuPDF (fitz) nor pypdf is installed. "
            "Install with: pip install pymupdf>=1.24.0 or pypdf>=4.0.0"
        )
        return []


class JSONLLoader(BaseLoader):
    """Loader for newline-delimited JSON files (.jsonl, .ndjson). Each line becomes a document."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        stable_base = _stable_file_id(path, "jsonl")
        documents = []
        with open(path, encoding="utf-8") as f:
            for i, line in enumerate(f):
                line = line.strip()
                if not line:
                    continue
                try:
                    item = json.loads(line)
                except json.JSONDecodeError:
                    logger.warning(f"Skipping malformed JSON at line {i + 1} in {path}")
                    continue
                if isinstance(item, dict):
                    text = item.get("text", item.get("content", json.dumps(item)))
                else:
                    text = json.dumps(item)
                documents.append(
                    {
                        "id": f"{stable_base}_{i}",
                        "text": text,
                        "metadata": {"source": path, "type": "jsonl", "line": i},
                    }
                )
        return documents


class NotebookLoader(BaseLoader):
    """Loader for Jupyter notebooks (.ipynb). Extracts markdown and code cells as documents."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        with open(path, encoding="utf-8") as f:
            try:
                nb = json.load(f)
            except json.JSONDecodeError as exc:
                logger.warning(f"Skipping malformed notebook {path}: {exc}")
                return []

        stable_base = _stable_file_id(path, "notebook")
        documents = []
        for i, cell in enumerate(nb.get("cells", [])):
            cell_type = cell.get("cell_type", "")
            source = cell.get("source", [])
            text = "".join(source) if isinstance(source, list) else source
            text = text.strip()
            if not text:
                continue
            if cell_type == "code":
                text = f"```python\n{text}\n```"
            documents.append(
                {
                    "id": f"{stable_base}_cell_{i}",
                    "text": text,
                    "metadata": {
                        "source": path,
                        "type": "notebook",
                        "cell_type": cell_type,
                        "cell_index": i,
                    },
                }
            )
        return documents


class ExcelLoader(BaseLoader):
    """Loader for Excel files (.xlsx, .xls) using pandas. Each sheet is loaded as a table."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        try:
            import pandas as pd
        except ImportError:
            logger.error("pandas not installed.")
            return []
        try:
            xl = pd.ExcelFile(path)
        except ImportError:
            logger.error("openpyxl not installed. Install with: pip install openpyxl")
            return []
        except Exception as exc:
            logger.warning(f"Could not open Excel file {path}: {exc}")
            return []

        from axon.splitters import TableSplitter

        stable_base = _stable_file_id(path, "excel")
        documents = []
        doc_idx = 0
        for sheet_name in xl.sheet_names:
            try:
                df = xl.parse(sheet_name).fillna("").astype(str)
            except Exception as exc:
                logger.warning(f"Skipping sheet '{sheet_name}' in {path}: {exc}")
                continue
            if df.empty:
                continue
            headers = list(df.columns)
            rows = df.to_dict(orient="records")
            splitter = TableSplitter(table_name=f"{os.path.basename(path)}[{sheet_name}]")
            for text in splitter.transform_rows(rows, headers):
                documents.append(
                    {
                        "id": f"{stable_base}_{sheet_name}_{doc_idx}",
                        "text": text,
                        "metadata": {
                            "source": path,
                            "type": "excel",
                            "sheet": sheet_name,
                            "row": doc_idx,
                        },
                    }
                )
                doc_idx += 1
        return documents


class ParquetLoader(BaseLoader):
    """Loader for Parquet files using pandas + pyarrow."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        try:
            import pandas as pd
        except ImportError:
            logger.error("pandas not installed.")
            return []
        try:
            df = pd.read_parquet(path).fillna("").astype(str)
        except ImportError:
            logger.error("pyarrow not installed. Install with: pip install pyarrow")
            return []
        except Exception as exc:
            logger.warning(f"Could not read Parquet file {path}: {exc}")
            return []

        from axon.splitters import TableSplitter

        stable_base = _stable_file_id(path, "parquet")
        splitter = TableSplitter(table_name=os.path.basename(path))
        chunks = splitter.transform_rows(df.to_dict(orient="records"), list(df.columns))
        return [
            {
                "id": f"{stable_base}_{i}",
                "text": text,
                "metadata": {"source": path, "type": "parquet", "row": i},
            }
            for i, text in enumerate(chunks)
        ]


class EPUBLoader(BaseLoader):
    """Loader for EPUB ebooks using ebooklib. Each document item becomes a chunk."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        try:
            import ebooklib
            from ebooklib import epub
        except ImportError:
            logger.error("ebooklib not installed. Install with: pip install ebooklib")
            return []
        try:
            book = epub.read_epub(path, options={"ignore_ncx": True})
        except Exception as exc:
            logger.warning(f"Could not read EPUB {path}: {exc}")
            return []

        stable_base = _stable_file_id(path, "epub")
        documents = []
        for i, item in enumerate(book.get_items_of_type(ebooklib.ITEM_DOCUMENT)):
            html = item.get_content().decode("utf-8", errors="ignore")
            text = _extract_html_text(html).strip()
            if not text:
                continue
            documents.append(
                {
                    "id": f"{stable_base}_chapter_{i}",
                    "text": text,
                    "metadata": {
                        "source": path,
                        "type": "epub",
                        "chapter": i,
                        "item_name": item.get_name(),
                    },
                }
            )
        return documents


class RTFLoader(BaseLoader):
    """Loader for RTF files using striprtf."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            logger.error("striprtf not installed. Install with: pip install striprtf")
            return []
        try:
            with open(path, encoding="utf-8", errors="ignore") as f:
                content = f.read()
            text = rtf_to_text(content).strip()
        except Exception as exc:
            logger.warning(f"Could not parse RTF file {path}: {exc}")
            return []
        return [
            {
                "id": _stable_file_id(path, "rtf"),
                "text": text,
                "metadata": {"source": path, "type": "rtf"},
            }
        ]


class XMLLoader(BaseLoader):
    """Loader for XML files. Strips tags and extracts clean text content."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        import xml.etree.ElementTree as ET

        try:
            tree = ET.parse(path)
        except ET.ParseError as exc:
            logger.warning(f"Could not parse XML {path}: {exc}")
            # Fall back to stripping tags with regex
            with open(path, encoding="utf-8", errors="ignore") as f:
                raw = f.read()
            text = re.sub(r"<[^>]+>", " ", raw)
            text = re.sub(r"\s+", " ", text).strip()
            return [
                {
                    "id": _stable_file_id(path, "xml"),
                    "text": text,
                    "metadata": {"source": path, "type": "xml"},
                }
            ]

        root = tree.getroot()
        parts: list[str] = []

        def _walk(node: ET.Element, depth: int = 0) -> None:
            tag = node.tag.split("}")[-1] if "}" in node.tag else node.tag
            attrs = " ".join(f'{k.split("}")[-1]}="{v}"' for k, v in node.attrib.items())
            header = f"<{tag}{' ' + attrs if attrs else ''}>"
            text = (node.text or "").strip()
            if text:
                parts.append(f"{header} {text}")
            elif node.attrib:
                parts.append(header)
            for child in node:
                _walk(child, depth + 1)
            tail = (node.tail or "").strip()
            if tail:
                parts.append(tail)

        _walk(root)
        full_text = "\n".join(parts)
        return [
            {
                "id": _stable_file_id(path, "xml"),
                "text": full_text,
                "metadata": {"source": path, "type": "xml"},
            }
        ]


class SQLLoader(BaseLoader):
    """Loader for SQL files. Extracts statements with natural-language descriptions."""

    # Patterns that identify DDL/DML statement types
    _STMT_RE = re.compile(
        r"^\s*(CREATE\s+(?:TABLE|VIEW|INDEX|PROCEDURE|FUNCTION|TRIGGER)|"
        r"INSERT\s+INTO|UPDATE\s+\w+\s+SET|DELETE\s+FROM|ALTER\s+TABLE|"
        r"DROP\s+(?:TABLE|VIEW)|SELECT\b)",
        re.IGNORECASE | re.MULTILINE,
    )

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        with open(path, encoding="utf-8", errors="ignore") as f:
            content = f.read()

        # Split on statement boundaries (semicolons followed by whitespace/newline)
        raw_stmts = re.split(r";\s*\n", content)
        stable_base = _stable_file_id(path, "sql")
        documents = []

        for i, stmt in enumerate(raw_stmts):
            stmt = stmt.strip()
            if not stmt or stmt.startswith("--"):
                continue

            # Extract comments preceding the statement as context
            comment_lines = [
                line.lstrip("- ").strip()
                for line in stmt.splitlines()
                if line.strip().startswith("--")
            ]
            code = "\n".join(
                line for line in stmt.splitlines() if not line.strip().startswith("--")
            ).strip()
            if not code:
                continue

            # Build a natural-language header from the first keyword + object name
            m = self._STMT_RE.match(code)
            if m:
                first_line = code.splitlines()[0]
                header = f"SQL statement: {first_line.strip()}"
            else:
                header = "SQL statement"

            comment_text = " ".join(comment_lines)
            text = f"{header}\n{comment_text}\n\n{code}".strip()

            documents.append(
                {
                    "id": f"{stable_base}_stmt_{i}",
                    "text": text,
                    "metadata": {"source": path, "type": "sql", "statement_index": i},
                }
            )

        # If no semicolon-delimited statements found, return the whole file
        if not documents:
            return [
                {
                    "id": stable_base,
                    "text": content.strip(),
                    "metadata": {"source": path, "type": "sql"},
                }
            ]

        return documents


class EMLLoader(BaseLoader):
    """Loader for .eml email files using stdlib email module."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        import email as _email
        import email.policy

        with open(path, "rb") as f:
            msg = _email.message_from_binary_file(f, policy=email.policy.default)

        from_addr = str(msg.get("From", ""))
        to_addr = str(msg.get("To", ""))
        subject = str(msg.get("Subject", ""))
        date = str(msg.get("Date", ""))

        # Extract body: prefer text/plain, fallback to text/html
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == "text/plain" and not body:
                    payload = part.get_payload(decode=True)
                    if isinstance(payload, bytes):
                        body = payload.decode(
                            part.get_content_charset() or "utf-8", errors="ignore"
                        )
                elif ct == "text/html" and not body:
                    payload = part.get_payload(decode=True)
                    if isinstance(payload, bytes):
                        html = payload.decode(
                            part.get_content_charset() or "utf-8", errors="ignore"
                        )
                        body = _extract_html_text(html)
        else:
            payload = msg.get_payload(decode=True)
            if isinstance(payload, bytes):
                body = payload.decode(msg.get_content_charset() or "utf-8", errors="ignore")
            if msg.get_content_type() == "text/html":
                body = _extract_html_text(body)

        text = (
            f"From: {from_addr}\nTo: {to_addr}\nSubject: {subject}\nDate: {date}\n\n{body}".strip()
        )
        return [
            {
                "id": _stable_file_id(path, "eml"),
                "text": text,
                "metadata": {
                    "source": path,
                    "type": "eml",
                    "email_from": from_addr,
                    "email_to": to_addr,
                    "email_subject": subject,
                    "email_date": date,
                },
            }
        ]


class MSGLoader(BaseLoader):
    """Loader for Outlook .msg files using extract-msg."""

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        try:
            import extract_msg
        except ImportError:
            logger.error("extract-msg not installed. Install with: pip install extract-msg")
            return []
        try:
            msg = extract_msg.Message(path)
        except Exception as exc:
            logger.warning(f"Could not open MSG file {path}: {exc}")
            return []

        from_addr = msg.sender or ""
        to_addr = msg.to or ""
        subject = msg.subject or ""
        date = str(msg.date) if msg.date else ""
        body = msg.body or ""
        if not body and msg.htmlBody:
            body = _extract_html_text(
                msg.htmlBody.decode("utf-8", errors="ignore")
                if isinstance(msg.htmlBody, bytes)
                else msg.htmlBody
            )

        text = (
            f"From: {from_addr}\nTo: {to_addr}\nSubject: {subject}\nDate: {date}\n\n{body}".strip()
        )
        return [
            {
                "id": _stable_file_id(path, "msg"),
                "text": text,
                "metadata": {
                    "source": path,
                    "type": "msg",
                    "email_from": from_addr,
                    "email_to": to_addr,
                    "email_subject": subject,
                    "email_date": date,
                },
            }
        ]


class LaTeXLoader(BaseLoader):
    """Loader for LaTeX .tex files. Strips commands and extracts prose content.

    Note: Mathematical equation content is intentionally discarded as it is
    not natural-language searchable. Use this loader for prose-heavy papers.
    """

    # Environments whose entire content is discarded (math, floats, bibliographies)
    _DISCARD_ENVS = re.compile(
        r"\\begin\{(figure|table|equation|align|align\*|eqnarray|eqnarray\*"
        r"|math|displaymath|tikzpicture|lstlisting|verbatim)\}.*?\\end\{\1\}",
        re.DOTALL,
    )
    # Commands that wrap prose — replace with their argument
    _WRAP_CMDS = re.compile(
        r"\\(?:textbf|textit|emph|texttt|underline|textrm|textsc|textsl"
        r"|section|subsection|subsubsection|paragraph|subparagraph"
        r"|chapter|title|author|date|caption|footnote)\*?\{([^}]*)\}",
        re.DOTALL,
    )
    # Non-prose commands to remove entirely (with optional bracket/brace args)
    _REMOVE_CMDS = re.compile(
        r"\\(?:label|ref|eqref|cite|citep|citet|citealt|nocite"
        r"|includegraphics|bibliography|bibliographystyle"
        r"|vspace|hspace|vskip|hskip|noindent|newpage|clearpage"
        r"|tableofcontents|listoffigures|listtables)"
        r"(?:\[[^\]]*\])?(?:\{[^}]*\})?",
    )
    # Remaining bare commands (e.g. \maketitle, \par, \LaTeX)
    _BARE_CMDS = re.compile(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?(?:\{[^}]*\})?")

    def load(self, path: str) -> list[dict[str, Any]]:
        _check_file_size(path)
        with open(path, encoding="utf-8", errors="ignore") as f:
            content = f.read()

        # Pass 1: strip comments (% to EOL, but not \%)
        content = re.sub(r"(?<!\\)%[^\n]*", "", content)
        # Pass 2: strip preamble (everything before \begin{document})
        doc_start = content.find(r"\begin{document}")
        if doc_start != -1:
            content = content[doc_start + len(r"\begin{document}") :]
        end_doc = content.find(r"\end{document}")
        if end_doc != -1:
            content = content[:end_doc]
        # Pass 3: discard math/float environments
        content = self._DISCARD_ENVS.sub(" ", content)
        # Pass 4: remove non-prose commands
        content = self._REMOVE_CMDS.sub(" ", content)
        # Pass 5: unwrap text-formatting commands
        content = self._WRAP_CMDS.sub(r"\1", content)
        # Pass 6: strip remaining bare commands
        content = self._BARE_CMDS.sub(" ", content)
        # Clean up braces, multiple spaces
        content = content.replace("{", " ").replace("}", " ")
        content = re.sub(r"\n{3,}", "\n\n", content)
        content = re.sub(r"[ \t]+", " ", content)
        text = content.strip()

        return [
            {
                "id": _stable_file_id(path, "latex"),
                "text": text,
                "metadata": {"source": path, "type": "latex"},
            }
        ]


class DirectoryLoader:
    """Loader that crawls a directory and uses appropriate loaders for each file."""

    def __init__(self, vlm_model: str = "llava"):
        _image_loader = ImageLoader(ollama_model=vlm_model)
        _excel_loader = ExcelLoader()
        _code_loader = CodeFileLoader()
        self.loaders = {
            # ── document / data formats ───────────────────────────────────
            ".txt": SmartTextLoader(),
            ".md": TextLoader(),
            ".tsv": FlexibleTableLoader(),
            ".json": JSONLoader(),
            ".jsonl": JSONLLoader(),
            ".ndjson": JSONLLoader(),
            ".ipynb": NotebookLoader(),
            ".csv": FlexibleTableLoader(),
            ".xlsx": _excel_loader,
            ".xls": _excel_loader,
            ".parquet": ParquetLoader(),
            ".epub": EPUBLoader(),
            ".rtf": RTFLoader(),
            ".eml": EMLLoader(),
            ".msg": MSGLoader(),
            ".tex": LaTeXLoader(),
            ".html": HTMLLoader(),
            ".htm": HTMLLoader(),
            ".xml": XMLLoader(),
            ".sql": SQLLoader(),
            ".docx": DOCXLoader(),
            ".pptx": PPTXLoader(),
            ".bmp": _image_loader,
            ".png": _image_loader,
            ".jpg": _image_loader,
            ".jpeg": _image_loader,
            ".tif": _image_loader,
            ".tiff": _image_loader,
            ".pgm": _image_loader,
            ".pdf": PDFLoader(),
            # ── source code formats ───────────────────────────────────────
            # CodeFileLoader sets type="code" so path-prefix enrichment is
            # skipped, keeping raw text valid for AST-based chunking.
            ".py": _code_loader,
            ".go": _code_loader,
            ".rs": _code_loader,
            ".cpp": _code_loader,
            ".c": _code_loader,
            ".h": _code_loader,
            ".hpp": _code_loader,
            ".sh": _code_loader,
            ".bash": _code_loader,
            ".zsh": _code_loader,
            ".rb": _code_loader,
            ".pl": _code_loader,
            ".pm": _code_loader,
            ".jl": _code_loader,
            ".js": _code_loader,
            ".jsx": _code_loader,
            ".ts": _code_loader,
            ".tsx": _code_loader,
            ".java": _code_loader,
            ".kt": _code_loader,
            ".cs": _code_loader,
            ".php": _code_loader,
            ".scala": _code_loader,
            ".swift": _code_loader,
        }

    def load(self, directory: str) -> list[dict[str, Any]]:
        all_documents = []
        path = Path(directory)

        for file_path in path.rglob("*"):
            if file_path.suffix.lower() in self.loaders:
                loader = self.loaders[file_path.suffix.lower()]
                try:
                    docs = loader.load(str(file_path))
                    # Native Path Enrichment (Breadcrumbs) for textual docs.
                    # Excluded: tabular/image types (structured data) and code
                    # files (type="code") — prepending into code text breaks
                    # AST-based chunkers like CodeAwareSplitter._split_python_ast.
                    for doc in docs:
                        if doc["metadata"].get("type") not in (
                            "csv",
                            "tsv",
                            "excel",
                            "parquet",
                            "image",
                            "code",
                        ):
                            rel_path = os.path.relpath(str(file_path), directory)
                            doc["text"] = f"[File Path: {rel_path}]\n{doc['text']}"
                    all_documents.extend(docs)
                except Exception as e:
                    logger.error(f"Failed to load {file_path}: {e}")

        return all_documents

    async def aload(self, directory: str) -> list[dict[str, Any]]:
        """Async version of load_directory — caps concurrency at 32 tasks."""
        path = Path(directory)
        semaphore = asyncio.Semaphore(32)

        async def _load_with_semaphore(loader, file_path: str):
            async with semaphore:
                return await loader.aload(file_path)

        tasks = []
        file_paths_for_tasks = []
        for file_path in path.rglob("*"):
            suffix = file_path.suffix.lower()
            if suffix in self.loaders:
                loader = self.loaders[suffix]
                tasks.append(_load_with_semaphore(loader, str(file_path)))
                file_paths_for_tasks.append(file_path)

        if not tasks:
            return []

        results = await asyncio.gather(*tasks, return_exceptions=True)
        all_documents = []
        for file_path, res in zip(file_paths_for_tasks, results):
            if isinstance(res, BaseException):
                if isinstance(res, asyncio.CancelledError):
                    raise res
                logger.warning("async file load failed: %s", res)
                continue
            for doc in res:
                if doc["metadata"].get("type") not in ("csv", "tsv", "image", "code"):
                    rel_path = os.path.relpath(str(file_path), directory)
                    doc["text"] = f"[File Path: {rel_path}]\n{doc['text']}"
            all_documents.extend(res)

        return all_documents
